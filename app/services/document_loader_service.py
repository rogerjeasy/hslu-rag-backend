# app/services/document_loader_service.py
import os
import logging
from typing import List, Dict, Any, Optional

logger = logging.getLogger(__name__)

class DocumentLoaderService:
    """Service for loading and extracting text from different file types"""
    
    def __init__(self):
        """Initialize document loader service"""
        pass
        
    async def extract_text(self, file_path: str, file_type: str) -> str:
        """
        Extract text from a file based on its type.
        
        Args:
            file_path: Path to the file
            file_type: Type of file (pdf, docx, etc.)
            
        Returns:
            Extracted text content
        """
        # Choose the appropriate loader based on file type
        if file_type == "pdf":
            return await self._extract_from_pdf(file_path)
        elif file_type in ["docx", "doc"]:
            return await self._extract_from_docx(file_path)
        elif file_type in ["pptx", "ppt"]:
            return await self._extract_from_ppt(file_path)
        elif file_type in ["xlsx", "xls"]:
            return await self._extract_from_excel(file_path)
        elif file_type == "csv":
            return await self._extract_from_csv(file_path)
        elif file_type in ["txt", "py", "js", "html", "css", "md"]:
            return await self._extract_from_text(file_path)
        elif file_type == "ipynb":
            return await self._extract_from_notebook(file_path)
        else:
            # Default to text extraction for unknown types
            return await self._extract_from_text(file_path)
    
    async def _extract_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF files"""
        try:
            # Try using PyPDF first (more reliable)
            import pypdf
            text_content = ""
            
            try:
                with open(file_path, 'rb') as file:
                    pdf = pypdf.PdfReader(file)
                    num_pages = len(pdf.pages)
                    
                    for i in range(num_pages):
                        page = pdf.pages[i]
                        page_text = page.extract_text() or ""
                        text_content += f"\n--- Page {i+1} ---\n{page_text}\n\n"
                        
                if text_content.strip():
                    return text_content
            except Exception as e:
                logger.warning(f"PyPDF extraction failed: {str(e)}, falling back to LangChain")
            
            # Fallback to LangChain if PyPDF fails or returns empty text
            from langchain.document_loaders import PyPDFLoader
            loader = PyPDFLoader(file_path)
            documents = loader.load()
            
            text_content = ""
            for i, doc in enumerate(documents):
                page_info = f"\n--- Page {i+1} ---\n" if hasattr(doc, 'metadata') and 'page' in doc.metadata else ""
                text_content += page_info + doc.page_content + "\n\n"
                
            return text_content
            
        except Exception as e:
            logger.error(f"Error extracting text from PDF {file_path}: {str(e)}")
            raise ValueError(f"Error loading {file_path}: {str(e)}")
    
    async def _extract_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX files"""
        try:
            # Try using python-docx2txt first (more reliable for docx)
            import docx2txt
            
            try:
                text_content = docx2txt.process(file_path)
                if text_content.strip():
                    return text_content
            except Exception as e:
                logger.warning(f"docx2txt extraction failed: {str(e)}, falling back to python-docx")
            
            # Fallback to python-docx if docx2txt fails
            import docx
            
            try:
                doc = docx.Document(file_path)
                text_content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
                if text_content.strip():
                    return text_content
            except Exception as e:
                logger.warning(f"python-docx extraction failed: {str(e)}, falling back to LangChain")
            
            # Final fallback to LangChain
            from langchain.document_loaders import Docx2txtLoader
            loader = Docx2txtLoader(file_path)
            documents = loader.load()
            
            text_content = ""
            for doc in documents:
                text_content += doc.page_content + "\n\n"
                
            return text_content
            
        except Exception as e:
            logger.error(f"Error extracting text from DOCX {file_path}: {str(e)}")
            raise ValueError(f"Error loading {file_path}: {str(e)}")
    
    async def _extract_from_ppt(self, file_path: str) -> str:
        """Extract text from PPT/PPTX files"""
        try:
            # Try using python-pptx first
            import pptx
            
            try:
                presentation = pptx.Presentation(file_path)
                text_content = ""
                
                for i, slide in enumerate(presentation.slides):
                    slide_text = ""
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text += shape.text + "\n"
                    
                    text_content += f"\n--- Slide {i+1} ---\n{slide_text}\n\n"
                
                if text_content.strip():
                    return text_content
            except Exception as e:
                logger.warning(f"python-pptx extraction failed: {str(e)}, falling back to LangChain")
            
            # Fallback to LangChain
            from langchain.document_loaders import UnstructuredPowerPointLoader
            loader = UnstructuredPowerPointLoader(file_path)
            documents = loader.load()
            
            text_content = ""
            for i, doc in enumerate(documents):
                slide_info = f"\n--- Slide {i+1} ---\n"
                text_content += slide_info + doc.page_content + "\n\n"
                
            return text_content
            
        except Exception as e:
            logger.error(f"Error extracting text from PPT {file_path}: {str(e)}")
            raise ValueError(f"Error loading {file_path}: {str(e)}")
    
    async def _extract_from_excel(self, file_path: str) -> str:
        """Extract text from Excel files"""
        try:
            # Try using pandas first
            import pandas as pd
            
            try:
                # Read all sheets
                xlsx = pd.ExcelFile(file_path)
                text_content = ""
                
                for sheet_name in xlsx.sheet_names:
                    df = pd.read_excel(xlsx, sheet_name)
                    sheet_text = f"\n--- Sheet: {sheet_name} ---\n"
                    sheet_text += df.to_string(index=False) + "\n\n"
                    text_content += sheet_text
                
                if text_content.strip():
                    return text_content
            except Exception as e:
                logger.warning(f"pandas Excel extraction failed: {str(e)}, falling back to LangChain")
            
            # Fallback to LangChain
            from langchain.document_loaders import UnstructuredExcelLoader
            loader = UnstructuredExcelLoader(file_path)
            documents = loader.load()
            
            text_content = ""
            for doc in documents:
                text_content += doc.page_content + "\n\n"
                
            return text_content
            
        except Exception as e:
            logger.error(f"Error extracting text from Excel {file_path}: {str(e)}")
            raise ValueError(f"Error loading {file_path}: {str(e)}")
    
    async def _extract_from_csv(self, file_path: str) -> str:
        """Extract text from CSV files"""
        try:
            # Use pandas for CSV files
            import pandas as pd
            
            df = pd.read_csv(file_path)
            text_content = df.to_string(index=False)
            
            return text_content
            
        except Exception as e:
            logger.error(f"Error extracting text from CSV {file_path}: {str(e)}")
            raise ValueError(f"Error loading {file_path}: {str(e)}")
    
    async def _extract_from_text(self, file_path: str) -> str:
        """Extract text from plain text files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            # Try different encoding if UTF-8 fails
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    return file.read()
            except Exception as e:
                logger.error(f"Error reading text file with latin-1 encoding: {str(e)}")
                raise ValueError(f"Error loading {file_path}: {str(e)}")
        except Exception as e:
            logger.error(f"Error extracting text from text file {file_path}: {str(e)}")
            raise ValueError(f"Error loading {file_path}: {str(e)}")
    
    async def _extract_from_notebook(self, file_path: str) -> str:
        """Extract text from Jupyter notebooks"""
        try:
            import json
            
            with open(file_path, 'r', encoding='utf-8') as file:
                notebook = json.load(file)
            
            text_content = ""
            cell_count = 1
            
            for cell in notebook.get('cells', []):
                cell_type = cell.get('cell_type', '')
                source = cell.get('source', [])
                
                if isinstance(source, list):
                    source = ''.join(source)
                
                if cell_type == 'markdown':
                    text_content += f"\n--- Markdown Cell {cell_count} ---\n{source}\n\n"
                elif cell_type == 'code':
                    text_content += f"\n--- Code Cell {cell_count} ---\n{source}\n"
                    if 'outputs' in cell:
                        outputs = []
                        for output in cell['outputs']:
                            if 'text' in output:
                                if isinstance(output['text'], list):
                                    outputs.append(''.join(output['text']))
                                else:
                                    outputs.append(output['text'])
                            elif 'data' in output and 'text/plain' in output['data']:
                                data = output['data']['text/plain']
                                if isinstance(data, list):
                                    outputs.append(''.join(data))
                                else:
                                    outputs.append(data)
                        
                        if outputs:
                            text_content += "\n--- Output ---\n" + "\n".join(outputs) + "\n"
                
                cell_count += 1
            
            return text_content
            
        except Exception as e:
            logger.error(f"Error extracting text from notebook {file_path}: {str(e)}")
            raise ValueError(f"Error loading {file_path}: {str(e)}")