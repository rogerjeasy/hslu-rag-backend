# app/rag/document_processor.py
import os
import logging
import tempfile
import asyncio
from typing import Dict, List, Any, Optional, Tuple
import PyPDF2
from PIL import Image
import pytesseract
import pandas as pd
import numpy as np
import docx
import pptx
import nbformat
from nbformat import reads as nb_reads
import chardet
from pptx import Presentation
import openpyxl
import json
import re
import httpx
from io import BytesIO

from app.core.config import settings
from app.services.cloudinary_service import CloudinaryService
from app.rag_new.chuncker import TextChunker
from app.rag_new.embdeddings import EmbeddingService

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """
    Service for processing different types of documents and extracting text
    """
    
    def __init__(
        self,
        cloudinary_service: CloudinaryService,
        chunker: TextChunker,
        embedding_service: EmbeddingService
    ):
        self.cloudinary_service = cloudinary_service
        self.chunker = chunker
        self.embedding_service = embedding_service
        
    async def process_file(self, file_url: str, file_type: str, material_id: str) -> Tuple[List[Dict[str, Any]], List[str]]:
        """
        Process a file from its URL and extract text content
        
        Args:
            file_url: URL of the file to process
            file_type: Type of the file (pdf, jpg, etc.)
            material_id: ID of the material
            
        Returns:
            Tuple of (chunks, vector_ids)
        """
        try:
            # Download the file
            content = await self._download_file(file_url)
            
            # Extract text based on file type
            text = await self._extract_text(content, file_type.lower())
            
            if not text:
                logger.warning(f"No text extracted from file: {file_url}")
                return [], []
            
            # Chunk the text
            chunks = self.chunker.chunk_text(text)
            
            # Get metadata
            metadata = {
                "material_id": material_id,
                "source_url": file_url,
                "file_type": file_type
            }
            
            # Create vector embeddings
            vector_ids = await self._create_embeddings(chunks, metadata)
            
            # Return chunked data and vector IDs
            return chunks, vector_ids
            
        except Exception as e:
            logger.error(f"Error processing file: {str(e)}")
            raise
    
    async def _download_file(self, file_url: str) -> bytes:
        """
        Download a file from URL
        
        Args:
            file_url: URL of the file
            
        Returns:
            File content as bytes
        """
        async with httpx.AsyncClient() as client:
            response = await client.get(file_url)
            response.raise_for_status()
            return response.content
    
    async def _extract_text(self, content: bytes, file_type: str) -> str:
        """
        Extract text from file content based on file type
        
        Args:
            content: File content as bytes
            file_type: Type of the file
            
        Returns:
            Extracted text
        """
        # Create a temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file_type}') as temp_file:
            temp_file.write(content)
            temp_path = temp_file.name
        
        try:
            # Process based on file type
            if file_type in ['pdf']:
                return await self._extract_text_from_pdf(temp_path)
            
            elif file_type in ['png', 'jpg', 'jpeg', 'gif']:
                return await self._extract_text_from_image(temp_path)
            
            elif file_type in ['py', 'js', 'ts', 'html', 'css', 'r']:
                return await self._extract_text_from_code_file(temp_path)
                
            elif file_type in ['ipynb']:
                return await self._extract_text_from_notebook(temp_path)
                
            elif file_type in ['pptx', 'ppt']:
                return await self._extract_text_from_presentation(temp_path)
                
            elif file_type in ['docx', 'doc']:
                return await self._extract_text_from_document(temp_path)
                
            elif file_type in ['xlsx', 'xls', 'csv']:
                return await self._extract_text_from_spreadsheet(temp_path, file_type)
                
            else:
                logger.warning(f"Unsupported file type: {file_type}")
                # Try to read as text file
                return await self._extract_text_from_code_file(temp_path)
                
        finally:
            # Clean up temp file
            if os.path.exists(temp_path):
                os.unlink(temp_path)
    
    async def _extract_text_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF file"""
        text = ""
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page_num in range(len(reader.pages)):
                page = reader.pages[page_num]
                text += f"\n--- Page {page_num + 1} ---\n"
                text += page.extract_text() or ""
        return text
    
    async def _extract_text_from_image(self, file_path: str) -> str:
        """Extract text from image using OCR"""
        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            return text
        except Exception as e:
            logger.error(f"Error extracting text from image: {str(e)}")
            return ""
    
    async def _extract_text_from_code_file(self, file_path: str) -> str:
        """Extract text from code file"""
        try:
            # Try to detect the encoding
            with open(file_path, 'rb') as file:
                content = file.read()
                result = chardet.detect(content)
                encoding = result['encoding']
            
            # Read with detected encoding
            with open(file_path, 'r', encoding=encoding) as file:
                return file.read()
        except Exception as e:
            logger.error(f"Error extracting text from code file: {str(e)}")
            # Fallback to utf-8
            try:
                with open(file_path, 'r', encoding='utf-8') as file:
                    return file.read()
            except:
                return ""
    
    async def _extract_text_from_notebook(self, file_path: str) -> str:
        """Extract text from Jupyter notebook"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                notebook = json.load(file)
            
            text = ""
            # Parse as notebook
            if 'cells' in notebook:
                for i, cell in enumerate(notebook['cells']):
                    cell_type = cell.get('cell_type', '')
                    source = cell.get('source', [])
                    
                    # Convert source to string if it's a list
                    if isinstance(source, list):
                        source = ''.join(source)
                    
                    text += f"\n--- Cell {i+1} ({cell_type}) ---\n"
                    text += source
            
            return text
        except Exception as e:
            logger.error(f"Error extracting text from notebook: {str(e)}")
            return ""
    
    async def _extract_text_from_presentation(self, file_path: str) -> str:
        """Extract text from PowerPoint presentation"""
        try:
            # Open presentation
            presentation = Presentation(file_path)
            
            text = ""
            # Extract text from slides
            for i, slide in enumerate(presentation.slides):
                text += f"\n--- Slide {i+1} ---\n"
                
                # Get title
                if slide.shapes.title:
                    text += f"Title: {slide.shapes.title.text}\n"
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += f"{shape.text}\n"
            
            return text
        except Exception as e:
            logger.error(f"Error extracting text from presentation: {str(e)}")
            return ""
    
    async def _extract_text_from_document(self, file_path: str) -> str:
        """Extract text from Word document"""
        try:
            doc = docx.Document(file_path)
            
            text = ""
            # Extract paragraphs
            for para in doc.paragraphs:
                text += para.text + "\n"
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " | "
                    text += "\n"
            
            return text
        except Exception as e:
            logger.error(f"Error extracting text from document: {str(e)}")
            return ""
    
    async def _extract_text_from_spreadsheet(self, file_path: str, file_type: str) -> str:
        """Extract text from Excel or CSV file"""
        try:
            if file_type == 'csv':
                # Read CSV
                df = pd.read_csv(file_path)
            else:
                # Read Excel
                df = pd.read_excel(file_path)
            
            # Convert to string representation
            text = f"Column headers: {', '.join(df.columns)}\n\n"
            
            # Get data sample (first 100 rows)
            sample_size = min(100, len(df))
            text += f"Data sample ({sample_size} rows):\n"
            text += df.head(sample_size).to_string()
            
            # Add basic statistics for numerical columns
            text += "\n\nNumerical statistics:\n"
            numeric_cols = df.select_dtypes(include=[np.number]).columns
            if len(numeric_cols) > 0:
                text += df[numeric_cols].describe().to_string()
            
            return text
        except Exception as e:
            logger.error(f"Error extracting text from spreadsheet: {str(e)}")
            return ""
    
    async def _create_embeddings(self, chunks: List[Dict[str, Any]], metadata: Dict[str, Any]) -> List[str]:
        """
        Create embeddings for text chunks and store them in vector database
        
        Args:
            chunks: List of text chunks
            metadata: Metadata for the chunks
            
        Returns:
            List of vector IDs
        """
        vector_ids = []
        
        for i, chunk in enumerate(chunks):
            chunk_id = f"{metadata['material_id']}_{i}"
            
            # Create combined metadata
            chunk_metadata = {
                **metadata,
                "chunk_index": i,
                "total_chunks": len(chunks),
                "chunk_text_preview": chunk["text"][:100] if len(chunk["text"]) > 100 else chunk["text"]
            }
            
            # Create embedding
            try:
                vector_id = await self.embedding_service.create_embedding(
                    text=chunk["text"],
                    vector_id=chunk_id,
                    metadata=chunk_metadata
                )
                vector_ids.append(vector_id)
            except Exception as e:
                logger.error(f"Error creating embedding for chunk {i}: {str(e)}")
        
        return vector_ids