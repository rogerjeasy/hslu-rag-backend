# app/utils/file_processors/pptx_processor.py
import logging
import io
from typing import Dict, Tuple, Any
import asyncio
import re

from app.core.exceptions import DocumentProcessingException

logger = logging.getLogger(__name__)

class PPTXProcessor:
    """
    Processes PowerPoint documents to extract text and metadata.
    
    This class handles PPTX-specific extraction, including
    slide-by-slide content and structural information.
    """
    
    async def process(self, file_content: bytes) -> Tuple[str, Dict[str, Any]]:
        """
        Process a PowerPoint file to extract text and metadata.
        
        Args:
            file_content: Binary content of the PowerPoint file
            
        Returns:
            Tuple of (extracted text, PowerPoint metadata)
        """
        try:
            # Run extraction in a thread pool since PowerPoint processing can be CPU-intensive
            loop = asyncio.get_running_loop()
            return await loop.run_in_executor(
                None, self._extract_pptx_content, file_content
            )
        except Exception as e:
            logger.error(f"Error processing PowerPoint: {str(e)}")
            raise DocumentProcessingException(f"Failed to process PowerPoint: {str(e)}")
    
    def _extract_pptx_content(self, file_content: bytes) -> Tuple[str, Dict[str, Any]]:
        """
        Extract content and metadata from a PowerPoint file.
        
        This method uses python-pptx for extraction.
        
        Args:
            file_content: Binary content of the PowerPoint file
            
        Returns:
            Tuple of (extracted text, PowerPoint metadata)
        """
        try:
            from pptx import Presentation
        except ImportError:
            logger.error("python-pptx package is not installed")
            raise DocumentProcessingException("python-pptx package is required but not installed")
        
        pptx_file = io.BytesIO(file_content)
        
        try:
            presentation = Presentation(pptx_file)
            
            # Extract metadata
            metadata = self._extract_metadata(presentation)
            
            # Extract text
            full_text = ""
            
            for i, slide in enumerate(presentation.slides):
                # Add slide marker
                full_text += f"\n\n--- Slide {i+1} ---\n\n"
                
                # Extract slide title if available
                if slide.shapes.title:
                    title_text = slide.shapes.title.text.strip()
                    if title_text:
                        full_text += f"Title: {title_text}\n\n"
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Skip if it's the title we already added
                        if shape == slide.shapes.title:
                            continue
                        full_text += f"{shape.text.strip()}\n"
                
                # Add notes if available
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    full_text += f"\nNotes: {notes_text}\n"
            
            return full_text, metadata
            
        except Exception as e:
            logger.error(f"Error extracting content from PowerPoint: {str(e)}")
            raise DocumentProcessingException(f"Failed to extract content from PowerPoint: {str(e)}")
    
    def _extract_metadata(self, presentation) -> Dict[str, Any]:
        """
        Extract metadata from a PowerPoint presentation.
        
        Args:
            presentation: python-pptx Presentation object
            
        Returns:
            Dictionary of metadata
        """
        metadata = {
            "num_slides": len(presentation.slides),
            "content_type": "slides"
        }
        
        # Extract document properties if available
        if hasattr(presentation, "core_properties"):
            props = presentation.core_properties
            if props.title:
                metadata["pptx_title"] = props.title
            if props.author:
                metadata["pptx_author"] = props.author
            if props.subject:
                metadata["pptx_subject"] = props.subject
            if props.keywords:
                metadata["pptx_keywords"] = props.keywords
            if props.created:
                metadata["pptx_created"] = props.created.isoformat()
            if props.modified:
                metadata["pptx_modified"] = props.modified.isoformat()
        
        return metadata